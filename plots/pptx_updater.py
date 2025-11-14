from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE # <-- AÑADIR

# --- AÑADIR 'imagenes_config=None' a la firma de la función ---
def actualizar_graficos(pptx_path, output_path, graficos_config, textos_config=None, imagenes_config=None):
    prs = Presentation(pptx_path)
    charts_map = {}
    text_shapes_map = {}
    image_placeholders = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            if isinstance(shape, GraphicFrame) and shape.has_chart:
                key = (shape.name or "").strip().lower()
                charts_map[key] = shape.chart
            
            elif shape.has_text_frame:
                key = (shape.name or "").strip().lower()
                text_shapes_map[key] = shape

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                key = (shape.name or "").strip().lower()
                if key:
                    image_placeholders[key] = (slide, shape)

    charts_no_porcentaje = [
        "nivel_bienestar",
        "nivel_bienestar_sede",
        "perma_2022_2025_chart",
        "nhl_2022_2025_chart"
    ]

    charts_no_porcentaje_lower = [name.lower() for name in charts_no_porcentaje]
    
    for chart_name, cfg in graficos_config.items():
        chart_name_lower = chart_name.lower()
        chart = charts_map.get(chart_name_lower)
        if not chart:
            continue

        df_preparado = cfg["dataset_fn"](*cfg["dataset_args"])
        chart_data = cfg["chart_builder"](df_preparado, *cfg["chart_builder_args"])
        chart.replace_data(chart_data)

        if chart_name_lower not in charts_no_porcentaje_lower:
            try:
                if hasattr(chart, "value_axis"):
                    chart.value_axis.tick_labels.number_format = "0%"
                elif hasattr(chart, "category_axis"):
                    chart.category_axis.tick_labels.number_format = "0%"
            except ValueError:
                pass

            for series in chart.series:
                series.has_data_labels = True
                series.data_labels.number_format = '0%'
        else:
            try:
                if hasattr(chart, "value_axis"):
                    chart.value_axis.tick_labels.number_format = "0.0" 
            except ValueError:
                pass
            
            for series in chart.series:
                series.has_data_labels = True
                series.data_labels.number_format = '0.0'

    if textos_config:
        for shape_name, cfg_texto in textos_config.items():
            shape = text_shapes_map.get(shape_name.lower())
            if not shape:
                continue

            if isinstance(cfg_texto, str):
                nuevo_texto = cfg_texto
                estilo = {}
            else:
                nuevo_texto = cfg_texto.get("texto", "")
                estilo = cfg_texto.get("estilo", {})

            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(nuevo_texto)

            font = run.font
            if "fuente" in estilo:
                font.name = estilo["fuente"]
            if "tamano" in estilo:
                font.size = Pt(estilo["tamano"])
            if "bold" in estilo:
                font.bold = estilo["bold"]
            if "italic" in estilo:
                font.italic = estilo["italic"]
            if "color" in estilo:
                r, g, b = estilo["color"]
                font.color.rgb = RGBColor(r, g, b)

    if imagenes_config:
        for shape_name, cfg_img in imagenes_config.items():
            key = shape_name.lower()
            if key not in image_placeholders:
                continue
            
            buffer = cfg_img.get("buffer")
            if not buffer or buffer.getbuffer().nbytes == 0:
                continue

            slide, placeholder = image_placeholders[key]
            
            left, top, width, height = (
                placeholder.left, placeholder.top, 
                placeholder.width, placeholder.height
            )

            try:
                slide.shapes.add_picture(buffer, left, top, width, height)

                sp = placeholder._element
                sp.getparent().remove(sp)
                
            
            except Exception as e:
                print(f"❌ Error al actualizar imagen '{shape_name}': {e}")

    prs.save(output_path)